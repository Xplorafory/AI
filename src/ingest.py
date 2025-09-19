# src/ingest.py
import os, sys, json, uuid, glob, re
from pathlib import Path
from datetime import datetime
from typing import List, Dict, Optional

import psycopg2
from psycopg2.extras import execute_batch

# --- Env loading (supports both .env and infra\infra.env)
from dotenv import load_dotenv, find_dotenv
found = find_dotenv(filename=".env", usecwd=True)
if found:
    load_dotenv(found)
repo_root = Path(__file__).resolve().parents[1]
infra_env = repo_root / "infra" / "infra.env"
if infra_env.exists():
    load_dotenv(infra_env.as_posix(), override=True)

# --- OpenAI / local embeddings
from openai import OpenAI
import tiktoken

# Optional parsers
from pypdf import PdfReader
from pptx import Presentation
import srt, webvtt

# Optional extras (docx + image OCR)
DOCX_AVAILABLE = False
try:
    from docx import Document
    DOCX_AVAILABLE = True
except Exception:
    pass

OCR_AVAILABLE = False
try:
    from PIL import Image
    import pytesseract
    OCR_AVAILABLE = True
except Exception:
    pass

# ----------------- Config -----------------
EMB_MODEL = os.getenv("EMB_MODEL", "text-embedding-3-small").strip()
PG_DSN    = os.getenv("PG_DSN", "postgresql://xploraforys:xploraforys@localhost:5432/xploraforys")
RAW_DIR   = os.getenv("RAW_DIR", "kb/raw")

# Map model -> embedding dimension
MODEL_DIMS = {
    "text-embedding-3-small": 1536,
    "text-embedding-3-large": 3072,
    # local fallback
    "local-MiniLM-384": 384,
}
# Table naming per dim
def table_for_dim(dim: int) -> str:
    return f"kb_chunks_{dim}"

# Clean up API key: strip whitespace/newlines
def clean_api_key(k: Optional[str]) -> Optional[str]:
    if not k: return None
    k = k.strip()
    # forbid spaces or newlines inside key
    if re.search(r"\s", k):
        print("❌ OPENAI_API_KEY contains whitespace/newlines. Fix your env variable.")
        return None
    return k

OPENAI_API_KEY = clean_api_key(os.getenv("OPENAI_API_KEY"))
use_openai = OPENAI_API_KEY is not None

print(f"🔐 OPENAI_API_KEY: {'set' if use_openai else 'not set'}")
print(f"🗄️  PG_DSN       : {PG_DSN}")
print(f"🧠 EMB_MODEL    : {EMB_MODEL}")

client = None
if use_openai:
    client = OpenAI(api_key=OPENAI_API_KEY)

# Tokenizer (used for chunking)
enc = tiktoken.get_encoding("cl100k_base")

# ----------------- Loaders -----------------
def load_txt(p):  return open(p, "r", encoding="utf-8", errors="ignore").read()
def load_md(p):   return load_txt(p)

def load_pdf(p):
    reader = PdfReader(p)
    return "\n".join([(pg.extract_text() or "") for pg in reader.pages])

def load_pptx(p):
    prs = Presentation(p)
    out = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                out.append(shape.text)
    return "\n".join(out)

def load_srt_file(p):
    with open(p, "r", encoding="utf-8", errors="ignore") as f:
        subs = list(srt.parse(f.read()))
    return "\n".join([s.content for s in subs])

def load_vtt_file(p):
    return "\n".join([c.text for c in webvtt.read(p)])

def load_docx(p):
    if not DOCX_AVAILABLE:
        print(f"⚠️ Skipping DOCX (missing python-docx): {p}")
        return ""
    doc = Document(p)
    return "\n".join([para.text for para in doc.paragraphs])

def load_image_ocr(p):
    if not OCR_AVAILABLE:
        print(f"⚠️ Skipping image OCR (pytesseract/Pillow missing): {p}")
        return ""
    try:
        img = Image.open(p)
        text = pytesseract.image_to_string(img, lang="eng")  # adjust language if you need 'swe'
        return text
    except Exception as e:
        print(f"⚠️ OCR error for {p}: {e}")
        return ""

# ----------------- Discover files -----------------
def discover_docs() -> List[Dict]:
    base = Path(RAW_DIR)
    paths = [p for p in base.rglob("*") if p.is_file()]
    docs = []
    for p in paths:
        pl = p.as_posix().lower()
        text = ""
        try:
            if pl.endswith((".txt", ".md")):
                text = load_txt(p) if pl.endswith(".txt") else load_md(p)
            elif pl.endswith(".pdf"):
                text = load_pdf(p)
            elif pl.endswith(".pptx"):
                text = load_pptx(p)
            elif pl.endswith(".srt"):
                text = load_srt_file(p)
            elif pl.endswith(".vtt"):
                text = load_vtt_file(p)
            elif pl.endswith(".docx"):
                text = load_docx(p)
            elif pl.endswith((".png", ".jpg", ".jpeg", ".tif", ".tiff")):
                text = load_image_ocr(p)
        except Exception as e:
            print(f"⚠️ Could not parse {p}: {e}")

        if text and text.strip():
            docs.append({"source": p.as_posix(), "title": p.name, "text": text})
    return docs

# ----------------- Chunking -----------------
CHUNK_TOKENS  = int(os.getenv("CHUNK_TOKENS", "450"))
CHUNK_OVERLAP = int(os.getenv("CHUNK_OVERLAP", "60"))

def chunk_text(text: str, max_tokens=CHUNK_TOKENS, overlap=CHUNK_OVERLAP):
    tokens = enc.encode(text)
    i = 0
    while i < len(tokens):
        chunk_tokens = tokens[i:i+max_tokens]
        yield enc.decode(chunk_tokens)
        i += max(1, max_tokens - overlap)

# ----------------- Tables / DB -----------------
def ensure_tables(conn, dim: int) -> str:
    tbl = table_for_dim(dim)
    with conn.cursor() as cur:
        cur.execute("CREATE EXTENSION IF NOT EXISTS vector;")
        cur.execute(f"""
            CREATE TABLE IF NOT EXISTS {tbl} (
              id TEXT PRIMARY KEY,
              content TEXT,
              embedding vector({dim}),
              source TEXT,
              title TEXT,
              created_at TIMESTAMP DEFAULT NOW()
            );
        """)
        cur.execute(f"""
            CREATE INDEX IF NOT EXISTS idx_{tbl}_embedding
              ON {tbl} USING ivfflat (embedding vector_cosine_ops) WITH (lists = 100);
        """)
    conn.commit()
    return tbl

# ----------------- Embeddings -----------------
def embed_batches_openai(texts: List[str], model: str, batch_size=64) -> List[List[float]]:
    out = []
    for i in range(0, len(texts), batch_size):
        batch = texts[i:i+batch_size]
        r = client.embeddings.create(model=model, input=batch)
        out.extend([d.embedding for d in r.data])
    return out

def embed_batches_local(texts: List[str]) -> List[List[float]]:
    # lightweight local embedding (MiniLM-384)
    from sentence_transformers import SentenceTransformer
    model = SentenceTransformer("sentence-transformers/all-MiniLM-L6-v2")
    return model.encode(texts, convert_to_numpy=False).tolist()

def model_dim(model_name: str) -> int:
    if use_openai:
        return MODEL_DIMS.get(model_name, 1536)
    return 384  # local

def upsert_rows(conn, table: str, rows: List[Dict], dim: int):
    sql = f"""
        INSERT INTO {table} (id, content, embedding, source, title, created_at)
        VALUES (%s, %s, %s, %s, %s, %s)
        ON CONFLICT (id) DO UPDATE
          SET content=EXCLUDED.content,
              embedding=EXCLUDED.embedding,
              source=EXCLUDED.source,
              title=EXCLUDED.title
    """
    vals = [(r["id"], r["content"], r["embedding"], r["source"], r["title"], datetime.utcnow()) for r in rows]
    with conn.cursor() as cur:
        execute_batch(cur, sql, vals, page_size=200)
    conn.commit()

# ----------------- Main -----------------
if __name__ == "__main__":
    dim = model_dim(EMB_MODEL)
    table = table_for_dim(dim)
    print(f"🗄️  Using table  : {table} (dim={dim})")

    docs = discover_docs()
    if not docs:
        print(f"⚠️ No documents found in {RAW_DIR}/")
        sys.exit(0)

    # build chunks
    chunks = []
    manifest = {}
    for d in docs:
        c_list = list(chunk_text(d["text"]))
        for c in c_list:
            chunks.append({
                "id": str(uuid.uuid4()),
                "content": c,
                "source": d["source"],
                "title": d["title"]
            })
        manifest.setdefault(d["title"], 0)
        manifest[d["title"]] += len(c_list)

    # embed
    contents = [c["content"] for c in chunks]
    if use_openai:
        print(f"✅ OpenAI embeddings model: {EMB_MODEL} ({dim} dim).")
        embs = embed_batches_openai(contents, EMB_MODEL, batch_size=64)
    else:
        print("ℹ️ No OPENAI_API_KEY set or invalid. Using local embeddings (MiniLM-384).")
        embs = embed_batches_local(contents)

    # attach
    for r, e in zip(chunks, embs):
        r["embedding"] = e

    # DB
    conn = psycopg2.connect(PG_DSN)
    tbl = ensure_tables(conn, dim)
    upsert_rows(conn, tbl, chunks, dim)
    print(f"✅ Upserted {len(chunks)} chunks → pgvector\n")

    # summary
    print("----- Ingestion summary -----")
    for k, v in manifest.items():
        print(f"{k}: {v} chunks")
    print(f"TOTAL: {sum(manifest.values())} chunks")
    print("-----------------------------")
